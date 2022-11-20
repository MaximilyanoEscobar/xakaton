from pptx import Presentation
from pptx.util import Inches
from PIL import Image

root = Presentation('Шаблон.pptx')
img_path = 'Gistogramma.png'
im = Image.open(img_path)
const = 0
width, height = im.size
max_height = 285
max_width = 427

while width > max_width or height > max_height:
    if width > max_width:
        const = max_width / width
        width = max_width
        height = height * const
    elif height > max_height:
        const = max_height / height
        height = max_height
        width = width * const

slide_one = root.slides[0]
for shape in slide_one.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = 'Нихера не продукт'


left = Inches(5.7)
top = Inches(4.1)
pic = slide_one.shapes.add_picture(img_path, left, top, width=Inches(width / 100), height=Inches(height / 100))
root.save('Шаблон тест.pptx')