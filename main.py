# -*- coding: utf-8 -*-
import json
import time
import datetime
import requests
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters import Text
from aiogram.dispatcher.filters.state import StatesGroup, State
from aiogram.types import CallbackQuery
from matplotlib import pyplot as plt
from Settings import BOT_TOKEN, MainMenu, AdminMenu, GoToQuestions, AdminStopKeyboard
from aiogram import Bot, Dispatcher, executor, types

bot = Bot(token=BOT_TOKEN, parse_mode="html")
storage = MemoryStorage()
dp = Dispatcher(bot, storage=storage)

with open("Questions.json", 'r') as file_questions:
    file_questions = json.load(file_questions)

with open("QuestionIndex.json", 'r') as file_questions_index:
    file_questions_index = json.load(file_questions_index)

with open('Users.json', 'r') as file_users:
    file_users = json.load(file_users)

with open("AdminList.json", 'r') as file_admins:
    file_admins = json.load(file_admins)


class InputMessage(StatesGroup):
    question = State()
    question_index_to_add = State()
    question_index_to_change = State()
    new_question = State()
    question_index_to_remove = State()
    question_to_remove_apply = State()
    questions_start = State()


@dp.message_handler(commands=['start'])
async def start(message: types.Message):
    await message.answer('<b>Приветствуем тебя в нашем боте!\nЗдесь ты cможешь отправить анкету со своими предложением.\n Анкетирование будет проходить в формате диалога с ботом!</b>', reply_markup=MainMenu)


@dp.message_handler(lambda message: message.text == "Заполнить заявку")
async def fill_application(message: types.Message):
    await message.answer("<b>Сейчас вам будет задан ряд вопросов. Вы готовы ответить на них?</b>", reply_markup=GoToQuestions)


@dp.message_handler(lambda message: message.text == "Справка")
async def fill_application(message: types.Message):
    await message.answer("Тут будет справка")


@dp.callback_query_handler(Text("Yes"))
async def start_questions(call: CallbackQuery):
    await bot.delete_message(call.from_user.id, call.message.message_id)
    if len(file_questions['Questions']) != 0:
        file_users[str(call.from_user.id)] = [[], {"TimeStart": None, "TimeEnd": None}]
        answers_list = []
        for id_question, question in enumerate(file_questions['Questions']):
            answers_list.append({"Question": question, "Answer": {"Icon": None, "Text": None}})
        file_users[str(call.from_user.id)][0] = answers_list
        file_users[str(call.from_user.id)][1]["TimeStart"] = int(time.time())
        with open("Users.json", 'w') as file_users_two:
            json.dump(file_users, file_users_two, indent=2)
        await call.message.answer('<b>Постарайтесь как можно точнее ответить на вопрос, нам будет интересно услышать вас</b>')
        await call.message.answer(f"<b>Вопрос №1</b>: <i>{file_questions['Questions'][0]}</i>")
        await InputMessage.questions_start.set()
    else:
        await call.message.answer('<b>Подождите, когда администрация добавит вопросы</b>')


@dp.message_handler(content_types=['photo', 'text'], state=InputMessage.questions_start)
async def send_question_to_user(message: types.Message, state: FSMContext):
    id_question = 0
    for id_question in range(len(file_users[str(message.from_user.id)][0])):
        if file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] is None and file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] is None:
            file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] = message.caption if message.caption is not None else message.text
            file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] = {"file_id": message.photo[-1]["file_id"], "file_unique_id": message.photo[-1]["file_unique_id"]} if len(message.photo) != 0 else None
            break
    if id_question == len(file_users[str(message.from_user.id)][0]) - 1:
        file_users[str(message.from_user.id)][1]['TimeEnd'] = int(time.time())
        with open("Users.json", 'w') as file_users_two:
            json.dump(file_users, file_users_two, indent=2)
        await message.answer('<b>Вы успешно прошли опрос. Ваша анкета сохранена и отправлена!</b>')
        await state.finish()
        root = Presentation('Шаблон.pptx')
        for index_question in range(7):
            slide = root.slides[index_question]
            if index_question == 0:
                if file_users[str(message.from_user.id)][0][index_question]["Answer"]["Icon"] is not None:
                    file_path = requests.get(f'https://api.telegram.org/bot{BOT_TOKEN}/getFile?file_id={file_users[str(message.from_user.id)][0][index_question]["Answer"]["Icon"]["file_id"]}').text
                    file_path = json.loads(file_path)
                    if file_path['result']:
                        file_path = str(file_path['result']['file_path'])
                        icon = requests.get(f'https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}').content
                        img_path = file_path.split("/")[1]
                        with open(img_path, 'wb') as new_file:
                            new_file.write(icon)
                        im = Image.open(img_path)
                        width, height = im.size
                        max_height = 285
                        max_width = 427
                        while width > max_width or height > max_height:
                            if width > max_width:
                                const = max_width / width
                                width = max_width
                                height = height * const
                            elif height > max_height:
                                const = max_height / height
                                height = max_height
                                width = width * const
                        left = Inches(5.5)
                        top = Inches(4.1)
                        slide.shapes.add_picture(img_path, left, top, width=Inches(width / 100), height=Inches(height / 100))
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
            elif index_question == 1:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
            elif index_question == 2:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
                if file_users[str(message.from_user.id)][0][index_question]["Answer"]["Icon"] is not None:
                    file_path = requests.get(
                        f'https://api.telegram.org/bot{BOT_TOKEN}/getFile?file_id={file_users[str(message.from_user.id)][0][index_question]["Answer"]["Icon"]["file_id"]}').text
                    file_path = json.loads(file_path)
                    if file_path['result']:
                        file_path = str(file_path['result']['file_path'])
                        icon = requests.get(f'https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}').content
                        img_path = file_path.split("/")[1]
                        with open(img_path, 'wb') as new_file:
                            new_file.write(icon)
                        im = Image.open(img_path)
                        width, height = im.size
                        max_height = 285
                        max_width = 427
                        while width > max_width or height > max_height:
                            if width > max_width:
                                const = max_width / width
                                width = max_width
                                height = height * const
                            elif height > max_height:
                                const = max_height / height
                                height = max_height
                                width = width * const
                        left = Inches(5.5)
                        top = Inches(0.15)
                        slide.shapes.add_picture(img_path, left, top, width=Inches(width / 100), height=Inches(height / 100))
            elif index_question == 3:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
            elif index_question == 4:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
            elif index_question == 5:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
            elif index_question == 6:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = file_users[str(message.from_user.id)][0][index_question]["Answer"]["Text"]
        root.save(f'{message.from_user.id}.pptx')
        with open(f'{message.from_user.id}.pptx', 'rb') as send_file:
            min_accepted = []
            for IdAdmin in file_admins:
                if file_admins[IdAdmin]['LastAccepted'] is None:
                    await bot.send_message(int(IdAdmin), '<b>Вам пришла новая анкета!</b>')
                    await bot.send_document(chat_id=int(IdAdmin), document=send_file)
                    file_admins[IdAdmin]['LastAccepted'] = int(time.time())
                    file_admins[IdAdmin]['Accepted'] += 1
                    with open('AdminList.json', 'w') as file_admins_two:
                        json.dump(file_admins, file_admins_two, indent=2)
                        return None
                else:
                    min_accepted.append(file_admins[IdAdmin]['Accepted'])
            min_accepted = min(min_accepted)
            min_accepted_admin_list = []
            for IdAdmin in file_admins:
                if file_admins[IdAdmin]['Accepted'] == min_accepted:
                    min_accepted_admin_list.append(IdAdmin)
            if len(min_accepted_admin_list) == 1:
                IdAdmin = min_accepted_admin_list[0]
                await bot.send_message(int(IdAdmin), '<b>Вам пришла новая анкета!</b>')
                await bot.send_document(chat_id=int(IdAdmin), document=send_file)
                file_admins[IdAdmin]['LastAccepted'] = int(time.time())
                file_admins[IdAdmin]['Accepted'] += 1
            else:
                min_timestamp = []
                for IdAdmin in min_accepted_admin_list:
                    min_timestamp.append(file_admins[IdAdmin]['LastAccepted'])
                min_timestamp = min(min_timestamp)
                for IdAdmin in file_admins:
                    if file_admins[IdAdmin]['LastAccepted'] == min_timestamp:
                        await bot.send_message(int(IdAdmin), '<b>Вам пришла новая анкета!</b>')
                        await bot.send_document(chat_id=int(IdAdmin), document=send_file)
                        file_admins[IdAdmin]['LastAccepted'] = int(time.time())
                        file_admins[IdAdmin]['Accepted'] += 1
                        break
        with open('AdminList.json', 'w') as file_admins_two:
            json.dump(file_admins, file_admins_two, indent=2)
    else:
        for id_question in range(len(file_users[str(message.from_user.id)][0])):
            if file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] is None and file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] is None:
                await message.answer(f"<b>Вопрос №{id_question + 1}:</b> <i>{file_users[str(message.from_user.id)][0][id_question]['Question']}</i>")
                break


@dp.callback_query_handler(Text(["No"]))
async def cansel_questions(call: CallbackQuery):
    await bot.delete_message(call.from_user.id, call.message.message_id)
    await call.message.answer('Выберите свои дальнейшие действия', reply_markup=MainMenu)


@dp.message_handler(commands=['admin'])
async def admin_settings(message: types.Message):
    if str(message.from_user.id) in file_admins:
        await message.answer('<b>Успешный вход в админ-панель!</b>', reply_markup=AdminMenu)
        requests_by_the_hour = [0 for _ in range(24)]
        timestamp = int(time.time())
        for TelegramId in file_users:
            if file_users[TelegramId][1]["TimeEnd"] is not None and datetime.datetime.fromtimestamp(file_users[TelegramId][1]["TimeEnd"]).strftime('%d') == datetime.datetime.fromtimestamp(timestamp).strftime('%d'):
                requests_by_the_hour[int(datetime.datetime.fromtimestamp(file_users[TelegramId][1]["TimeEnd"]).strftime('%H'))] += 1
        x_list = []
        for hour in range(24):
            hour = str(hour)
            if len(hour) == 1:
                hour = '0' + hour
            x_list.append(hour)
        y_list = requests_by_the_hour

        plt.title('Количество заявок по часам')
        plt.xlabel('Часы')
        plt.ylabel('Количество заявок')

        plt.bar(x_list, y_list)
        plt.savefig('Gistogramma.png', dpi='figure',
                    bbox_inches=None, pad_inches=1,
                    facecolor='auto', edgecolor='auto',
                    backend=None)
        with open('Gistogramma.png', 'rb') as file:
            await message.answer('<b>Статистика заявок за день:</b>')
            await bot.send_photo(message.from_user.id, photo=file)
    else:
        await message.answer('<b>Вы не являетесь администратором</b>!')


@dp.message_handler(lambda message: message.text == "Главное меню")
async def out_admin_settings(message: types.Message):
    await message.answer('Выберите свои дальнейшие действия', reply_markup=MainMenu)


@dp.message_handler(lambda message: message.text == "Список вопросов")
async def send_questions_to_admin(message: types.Message):
    if len(file_questions["Questions"]) != 0:
        questions = ''
        for id_question, Question in enumerate(file_questions["Questions"]):
            questions += f'<b>№{id_question + 1}:</b> <i>{Question}</i>\n\n'
        await message.answer(questions)
    else:
        await message.answer('<b>Список вопросов пуст</b>')


@dp.message_handler(lambda message: message.text == "Добавить вопрос")
async def send_new_question_out_admin(message: types.Message):
    await message.answer('<b>Каким по номеру должен стоять новый вопрос из всех вопросов:</b>',
                         reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_add.set()


@dp.message_handler(state=InputMessage.question_index_to_add)
async def add_new_question(message: types.Message, state: FSMContext):
    try:
        if message.text == 'Отменить действие':
            await state.finish()
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        elif int(message.text) > 0:
            file_questions_index[str(message.from_user.id)] = int(message.text)
            with open('QuestionIndex.json', 'w') as file_questions_index_two:
                json.dump(file_questions_index, file_questions_index_two, indent=2)
            await state.finish()
            await message.answer('<b>Введите ваш вопрос:</b>', reply_markup=AdminStopKeyboard)
            await InputMessage.question.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    except ExceptionGroup:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.question)
async def replace_new_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
    else:
        id_question = file_questions_index[str(message.from_user.id)]
        questions = list(file_questions["Questions"])[:id_question - 1] + [message.text] + list(
            file_questions["Questions"])[id_question - 1:]
        file_questions["Questions"] = questions
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2)
        await message.answer('<b>Вопрос успешно добавлен в список!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()


@dp.message_handler(lambda message: message.text == "Изменить вопрос")
async def change_question_index_out_admin(message: types.Message):
    await message.answer('<b>Пришлите номер вопроса, какой вы хотите изменить</b>', reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_change.set()


@dp.message_handler(state=InputMessage.question_index_to_change)
async def change_question(message: types.Message, state: FSMContext):
    try:
        if message.text == 'Отменить действие':
            await state.finish()
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        elif 0 < int(message.text) <= len(file_questions['Questions']):
            file_questions_index[str(message.from_user.id)] = int(message.text)
            with open('QuestionIndex.json', 'w') as file_questions_index_two:
                json.dump(file_questions_index, file_questions_index_two, indent=2)
            await state.finish()
            await message.answer('Введите ваш новый вопрос', reply_markup=AdminStopKeyboard)
            await InputMessage.new_question.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    except Exception:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.new_question)
async def replace_new_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
    else:
        id_question = file_questions_index[str(message.from_user.id)]
        file_questions["Questions"][id_question - 1] = message.text
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2)
        await message.answer('<b>Вопрос успешно изменён!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()


@dp.message_handler(lambda message: message.text == "Удалить вопрос")
async def send_question_index_to_remove_out_admin(message: types.Message):
    await message.answer('<b>Пришлите номер вопроса, какой вы хотите удалить</b>', reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_remove.set()


@dp.message_handler(state=InputMessage.question_index_to_remove)
async def remove_question_apply(message: types.Message, state: FSMContext):
    try:
        if message.text == 'Отменить действие':
            await state.finish()
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        elif 0 < int(message.text) <= len(file_questions['Questions']):
            file_questions_index[str(message.from_user.id)] = int(message.text)
            with open('QuestionIndex.json', 'w') as file_questions_index_two:
                json.dump(file_questions_index, file_questions_index_two, indent=2)
            await state.finish()
            await message.answer(
                f'<b>Вы точно хотите удалить вопрос:</b> <i>{file_questions["Questions"][int(message.text) - 1][:35]}</i>..?',
                reply_markup=AdminStopKeyboard)
            await message.answer('<b>Ответьте: <i>Да / Нет</i></b>', reply_markup=AdminStopKeyboard)
            await InputMessage.question_to_remove_apply.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    except Exception:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.question_to_remove_apply)
async def remove_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие' or message.text == 'Нет':
        await state.finish()
        await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
    elif message.text == 'Да':
        id_question = file_questions_index[str(message.from_user.id)]
        del (file_questions["Questions"][id_question - 1])
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2)
        await message.answer('<b>Вопрос успешно удалён!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()
    else:
        await message.answer('<b>Ответьте: <i>Да / Нет</i></b>', reply_markup=AdminStopKeyboard)


@dp.message_handler(content_types=["photo", "text"])
async def download_photo(message: types.Message, state: FSMContext):
    await send_question_to_user(message, state)

if __name__ == "__main__":
    executor.start_polling(dp, skip_updates=True, timeout=None)
